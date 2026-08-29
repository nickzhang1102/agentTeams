"""
文档解析模块
支持 PDF、Word、Excel、PowerPoint 格式的文本提取
"""

import io
import logging
from typing import Dict, Optional

logger = logging.getLogger(__name__)

# 文档解析资源上限（防止资源耗尽型 DoS）
MAX_PDF_PAGES = 500          # PDF 最大页数
MAX_DOCX_PARAGRAPHS = 50000  # Word 最大段落数
MAX_PPTX_SLIDES = 500        # PowerPoint 最大幻灯片数
MAX_EXCEL_ROWS = 100000      # Excel 单表最大行数（read_only 流式时限制）


class DocumentProcessor:
    """文档处理器，支持多种格式的文档解析"""

    @staticmethod
    def get_file_extension(filename: str) -> str:
        """获取文件扩展名（小写）"""
        if '.' not in filename:
            return ''
        return '.' + filename.rsplit('.', 1)[-1].lower()

    @staticmethod
    def is_supported(filename: str) -> bool:
        """检查文件是否支持解析"""
        ext = DocumentProcessor.get_file_extension(filename)
        supported = [
            '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'
        ]
        return ext in supported

    @staticmethod
    def process_document_from_path(filename: str, file_path: str) -> Dict:
        """
        从文件路径直接解析文档（流式读取，不将完整内容读入内存）

        各解析库原生支持文件路径参数，避免 bytes 全量加载。
        保留原 process_document(filename, bytes) 向后兼容。

        Args:
            filename: 文件名（用于扩展名判断）
            file_path: 文件在磁盘上的绝对路径

        Returns:
            dict: 同 process_document 返回格式
        """
        ext = DocumentProcessor.get_file_extension(filename)

        try:
            if ext == '.pdf':
                return DocumentProcessor._process_pdf_from_path(file_path)
            elif ext in ['.doc', '.docx']:
                return DocumentProcessor._process_word_from_path(file_path)
            elif ext in ['.xls', '.xlsx']:
                return DocumentProcessor._process_excel_from_path(file_path, filename)
            elif ext in ['.ppt', '.pptx']:
                return DocumentProcessor._process_ppt_from_path(file_path, filename)
            else:
                return {
                    'success': False,
                    'error': f'不支持的文件类型: {ext}',
                    'text': None
                }
        except Exception as e:
            logger.error(f"文档解析失败 [{filename}]: {str(e)}", exc_info=True)
            return {
                'success': False,
                'error': f'文档解析失败: {str(e)}',
                'text': None
            }

    # === 路径流式解析方法 ===

    @staticmethod
    def _process_pdf_from_path(file_path: str) -> Dict:
        """从文件路径解析 PDF（PyMuPDF 原生支持路径）"""
        try:
            import fitz  # pymupdf

            doc = fitz.open(file_path)
            try:

                page_count = len(doc)
                if page_count > MAX_PDF_PAGES:
                    return {
                        'success': False,
                        'error': f'PDF 页数超限（{page_count} > {MAX_PDF_PAGES}），请拆分或缩减后重试',
                        'text': None
                    }

                text_parts = []
                metadata = {'page_count': page_count, 'type': 'PDF 文档'}

                doc_metadata = doc.metadata
                if doc_metadata:
                    if doc_metadata.get('title'):
                        metadata['title'] = doc_metadata['title']
                    if doc_metadata.get('author'):
                        metadata['author'] = doc_metadata['author']
                    if doc_metadata.get('creator'):
                        metadata['creator'] = doc_metadata['creator']

                for i, page in enumerate(doc):
                    page_text = page.get_text("text")
                    if page_text.strip():
                        text_parts.append(f"--- 第 {i + 1} 页 ---\n{page_text}")

                full_text = '\n\n'.join(text_parts)

                if not full_text.strip():
                    return {
                        'success': True,
                        'text': '[此 PDF 可能是扫描版或图片型文档，无法提取文本内容。请考虑使用 OCR 工具处理。]',
                        'metadata': metadata,
                        'warning': '未提取到文本内容，可能是扫描版 PDF'
                    }

                return {'success': True, 'text': full_text, 'metadata': metadata}
            finally:
                doc.close()
        except ImportError:
            return {'success': False, 'error': 'PDF 解析库未安装，请安装 pymupdf', 'text': None}

    @staticmethod
    def _process_word_from_path(file_path: str) -> Dict:
        """从文件路径解析 Word（python-docx 原生支持路径）"""
        try:
            from docx import Document

            doc = Document(file_path)
            text_parts = []

            metadata = {}
            if doc.core_properties.title:
                metadata['title'] = doc.core_properties.title
            if doc.core_properties.author:
                metadata['author'] = doc.core_properties.author

            if len(doc.paragraphs) > MAX_DOCX_PARAGRAPHS:
                return {
                    'success': False,
                    'error': f'Word 段落数超限（{len(doc.paragraphs)} > {MAX_DOCX_PARAGRAPHS}）',
                    'text': None
                }

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        table_text.append(row_text)
                if table_text:
                    text_parts.append('\n[表格]\n' + '\n'.join(table_text))

            full_text = '\n\n'.join(text_parts)

            return {
                'success': True,
                'text': full_text,
                'metadata': {
                    **metadata,
                    'paragraph_count': len(doc.paragraphs),
                    'table_count': len(doc.tables),
                    'type': 'Word 文档'
                }
            }
        except ImportError:
            return {'success': False, 'error': 'Word 解析库未安装，请安装 python-docx', 'text': None}

    @staticmethod
    def _process_excel_from_path(file_path: str, filename: str) -> Dict:
        """从文件路径解析 Excel（openpyxl read_only 模式原生支持路径）"""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, data_only=True, read_only=True)
            try:
                text_parts = []

                metadata = {
                    'sheet_count': len(wb.sheetnames),
                    'sheet_names': wb.sheetnames,
                    'type': 'Excel 文档'
                }

                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    sheet_text = [f'=== 工作表: {sheet_name} ===']
                    row_limit_hit = False

                    for row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
                        if row_idx > MAX_EXCEL_ROWS:
                            row_limit_hit = True
                            break
                        row_data = [str(c) if c is not None else '' for c in row]
                        if any(row_data):
                            sheet_text.append(' | '.join(row_data))

                    if row_limit_hit:
                        sheet_text.append(f'... (已截断，行数超过 {MAX_EXCEL_ROWS})')
                    text_parts.append('\n'.join(sheet_text))

                full_text = '\n\n'.join(text_parts)

                return {'success': True, 'text': full_text, 'metadata': metadata}
            finally:
                wb.close()
        except ImportError:
            return {'success': False, 'error': 'Excel 解析库未安装，请安装 openpyxl', 'text': None}

    @staticmethod
    def _process_ppt_from_path(file_path: str, filename: str) -> Dict:
        """从文件路径解析 PowerPoint（python-pptx 原生支持路径）"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)

            slide_count = len(prs.slides)
            if slide_count > MAX_PPTX_SLIDES:
                return {
                    'success': False,
                    'error': f'PowerPoint 幻灯片数超限（{slide_count} > {MAX_PPTX_SLIDES}）',
                    'text': None
                }

            text_parts = []
            metadata = {'slide_count': slide_count, 'type': 'PowerPoint 文档'}

            for i, slide in enumerate(prs.slides, 1):
                slide_text = [f'=== 幻灯片 {i} ===']
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text)
                text_parts.append('\n'.join(slide_text))

            full_text = '\n\n'.join(text_parts)
            return {'success': True, 'text': full_text, 'metadata': metadata}
        except ImportError:
            return {'success': False, 'error': 'PowerPoint 解析库未安装，请安装 python-pptx', 'text': None}

    @staticmethod
    def process_document(filename: str, content: bytes) -> Dict:
        """
        处理文档并提取文本内容

        Args:
            filename: 文件名
            content: 文件二进制内容

        Returns:
            dict: 包含提取结果的字典
                - success: 是否成功
                - text: 提取的文本内容
                - error: 错误信息（如果失败）
                - metadata: 文档元数据
        """
        ext = DocumentProcessor.get_file_extension(filename)

        try:
            # 根据文件类型选择解析方法
            if ext == '.pdf':
                return DocumentProcessor._process_pdf(content)
            elif ext in ['.doc', '.docx']:
                return DocumentProcessor._process_word(content)
            elif ext in ['.xls', '.xlsx']:
                return DocumentProcessor._process_excel(content, filename)
            elif ext in ['.ppt', '.pptx']:
                return DocumentProcessor._process_ppt(content, filename)
            else:
                return {
                    'success': False,
                    'error': f'不支持的文件类型: {ext}',
                    'text': None
                }
        except Exception as e:
            logger.error(f"文档解析失败 [{filename}]: {str(e)}", exc_info=True)
            return {
                'success': False,
                'error': f'文档解析失败: {str(e)}',
                'text': None
            }

    @staticmethod
    def _process_pdf(content: bytes) -> Dict:
        """解析 PDF 文档（使用 PyMuPDF，更强大的文本提取能力）"""
        try:
            import fitz  # pymupdf

            # 从二进制内容打开文档
            doc = fitz.open(stream=content, filetype="pdf")
            try:

                # 页数上限检查，防止资源耗尽
                page_count = len(doc)
                if page_count > MAX_PDF_PAGES:
                    return {
                        'success': False,
                        'error': f'PDF 页数超限（{page_count} > {MAX_PDF_PAGES}），请拆分或缩减后重试',
                        'text': None
                    }

                text_parts = []

                # 提取元数据
                metadata = {
                    'page_count': page_count,
                    'type': 'PDF 文档'
                }

                # 提取文档元信息
                doc_metadata = doc.metadata
                if doc_metadata:
                    if doc_metadata.get('title'):
                        metadata['title'] = doc_metadata['title']
                    if doc_metadata.get('author'):
                        metadata['author'] = doc_metadata['author']
                    if doc_metadata.get('creator'):
                        metadata['creator'] = doc_metadata['creator']

                # 提取每一页的文本
                for i, page in enumerate(doc):
                    # 使用 get_text() 方法提取文本，支持多种输出格式
                    # "text" 是默认格式，保留文本布局
                    page_text = page.get_text("text")
                    if page_text.strip():
                        text_parts.append(f"--- 第 {i + 1} 页 ---\n{page_text}")

                full_text = '\n\n'.join(text_parts)

                # 关闭文档

                # 如果没有提取到任何文本，可能是扫描版 PDF
                if not full_text.strip():
                    return {
                        'success': True,
                        'text': '[此 PDF 可能是扫描版或图片型文档，无法提取文本内容。请考虑使用 OCR 工具处理。]',
                        'metadata': metadata,
                        'warning': '未提取到文本内容，可能是扫描版 PDF'
                    }

                return {
                    'success': True,
                    'text': full_text,
                    'metadata': metadata
                }
            finally:
                doc.close()
        except ImportError:
            return {
                'success': False,
                'error': 'PDF 解析库未安装，请安装 pymupdf',
                'text': None
            }

    @staticmethod
    def _process_word(content: bytes) -> Dict:
        """解析 Word 文档

        安全说明：python-docx 底层使用 lxml 解析 OOXML。lxml 默认禁用外部实体
        解析（resolve_entities=False 且禁用 DTD），故 XXE 攻击面已被库默认缓解。
        本方法额外增加段落数上限，防止资源耗尽。未使用 defusedxml 是因 python-docx
        不暴露可替换的 XML parser 接口。
        """
        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            text_parts = []

            # 提取核心属性
            metadata = {}
            if doc.core_properties.title:
                metadata['title'] = doc.core_properties.title
            if doc.core_properties.author:
                metadata['author'] = doc.core_properties.author

            # 段落数上限检查，防止资源耗尽
            if len(doc.paragraphs) > MAX_DOCX_PARAGRAPHS:
                return {
                    'success': False,
                    'error': f'Word 段落数超限（{len(doc.paragraphs)} > {MAX_DOCX_PARAGRAPHS}）',
                    'text': None
                }

            # 提取所有段落
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # 提取表格内容
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        table_text.append(row_text)
                if table_text:
                    text_parts.append('\n[表格]\n' + '\n'.join(table_text))

            full_text = '\n\n'.join(text_parts)

            return {
                'success': True,
                'text': full_text,
                'metadata': {
                    **metadata,
                    'paragraph_count': len(doc.paragraphs),
                    'table_count': len(doc.tables),
                    'type': 'Word 文档'
                }
            }
        except ImportError:
            return {
                'success': False,
                'error': 'Word 解析库未安装，请安装 python-docx',
                'text': None
            }

    @staticmethod
    def _process_excel(content: bytes, filename: str) -> Dict:
        """解析 Excel 文档

        使用 read_only=True 流式读取，降低大文件内存占用。read_only 模式下
        openpyxl 不缓存整表，缓解资源耗尽风险。
        """
        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
            try:
                text_parts = []

                # 提取元数据
                metadata = {
                    'sheet_count': len(wb.sheetnames),
                    'sheet_names': wb.sheetnames,
                    'type': 'Excel 文档'
                }

                # 提取每个工作表的内容（流式逐行读取）
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    sheet_text = [f'=== 工作表: {sheet_name} ===']
                    row_limit_hit = False

                    for row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
                        if row_idx > MAX_EXCEL_ROWS:
                            row_limit_hit = True
                            break
                        row_data = [str(c) if c is not None else '' for c in row]
                        if any(row_data):  # 只添加非空行
                            sheet_text.append(' | '.join(row_data))

                    if row_limit_hit:
                        sheet_text.append(f'... (已截断，行数超过 {MAX_EXCEL_ROWS})')
                    text_parts.append('\n'.join(sheet_text))

                full_text = '\n\n'.join(text_parts)

                # read_only 模式需显式关闭释放资源

                return {
                    'success': True,
                    'text': full_text,
                    'metadata': metadata
                }
            finally:
                wb.close()
        except ImportError:
            return {
                'success': False,
                'error': 'Excel 解析库未安装，请安装 openpyxl',
                'text': None
            }

    @staticmethod
    def _process_ppt(content: bytes, filename: str) -> Dict:
        """解析 PowerPoint 文档

        安全说明：python-pptx 底层使用 lxml 解析 OOXML，默认禁用外部实体解析，
        XXE 攻击面已由库默认缓解。本方法额外加幻灯片数上限防止资源耗尽。
        """
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))

            # 幻灯片数上限检查，防止资源耗尽
            slide_count = len(prs.slides)
            if slide_count > MAX_PPTX_SLIDES:
                return {
                    'success': False,
                    'error': f'PowerPoint 幻灯片数超限（{slide_count} > {MAX_PPTX_SLIDES}）',
                    'text': None
                }

            text_parts = []

            # 提取元数据
            metadata = {
                'slide_count': slide_count,
                'type': 'PowerPoint 文档'
            }

            # 提取每张幻灯片的文本
            for i, slide in enumerate(prs.slides, 1):
                slide_text = [f'=== 幻灯片 {i} ===']

                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text)

                text_parts.append('\n'.join(slide_text))

            full_text = '\n\n'.join(text_parts)

            return {
                'success': True,
                'text': full_text,
                'metadata': metadata
            }
        except ImportError:
            return {
                'success': False,
                'error': 'PowerPoint 解析库未安装，请安装 python-pptx',
                'text': None
            }


def process_uploaded_file(filename: str, content: bytes) -> Dict:
    """
    便捷函数：处理上传的文件

    Args:
        filename: 文件名
        content: 文件二进制内容

    Returns:
        dict: 处理结果
    """
    return DocumentProcessor.process_document(filename, content)